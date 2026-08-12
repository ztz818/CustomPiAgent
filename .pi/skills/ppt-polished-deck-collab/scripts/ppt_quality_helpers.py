#!/usr/bin/env python3
"""`ppt-polished-deck-collab` 的质量 gate 公共 helper。

定位与作用
----------
本文件服务三个 deck 级质量 gate：
1. `package_preflight`：检查 `pptx` 包内结构一致性与移动端兼容风险；
2. `structure_precheck`：检查文本框 fit、文字遮挡与结构化对象的排版边界风险。
3. `render_review`：检查预览图层面的边界触墨与扁平化图像风险。

它不直接做完整 gate 判断，而是提供统一的对象模型、PPT 遍历、几何换算、
文本归一化、问题汇总与 Markdown / JSON 输出能力，避免多个脚本重复发明接口。
"""

from __future__ import annotations

import json
import math
import re
from collections.abc import Iterable
from dataclasses import asdict, dataclass
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from agent_qc_reminders import build_agent_reminder
from agent_qc_reminders import quality_issues_to_observations
from agent_qc_reminders import sidecar_path
from agent_qc_reminders import write_agent_reminder

EMU_PER_PT = 12700.0
DEFAULT_BODY_FONT_PT = 14.0
FONT_SIZE_LOCATION_EXAMPLE_LIMIT = 12
TABLE_INDENT_TOLERANCE_PT = 0.25


@dataclass(frozen=True)
class RectPt:
    """以 pt 为单位的矩形。"""

    left: float
    top: float
    width: float
    height: float

    @property
    def right(self) -> float:
        """返回右边界。"""
        return self.left + self.width

    @property
    def bottom(self) -> float:
        """返回下边界。"""
        return self.top + self.height


@dataclass(frozen=True)
class ShapeRecord:
    """统一的 shape 级记录。"""

    slide_number: int
    shape_id: int
    shape_name: str
    shape_type: str
    z_order: int
    parent_group_id: int | None
    rect_pt: RectPt
    has_text: bool
    has_table: bool
    has_chart: bool
    is_picture: bool
    source_kind: str


@dataclass(frozen=True)
class TextItemRecord:
    """统一的文本承载对象记录。"""

    slide_number: int
    owner_shape_id: int
    owner_shape_name: str
    owner_shape_type: str
    z_order: int
    source_kind: str
    text: str
    rect_pt: RectPt
    inner_rect_pt: RectPt
    font_size_pt: float
    paragraph_count: int


@dataclass(frozen=True)
class FontSizeOccurrence:
    """单个可见文本 run 的显式字号记录。"""

    slide_number: int
    owner_shape_id: int
    owner_shape_name: str
    source_kind: str
    text: str
    font_size_pt: float
    setting_source: str


@dataclass(frozen=True)
class TableCellFormatRecord:
    """表格单元格级格式事实。"""

    slide_number: int
    owner_shape_id: int
    owner_shape_name: str
    table_index_on_slide: int
    row: int
    col: int
    vertical_anchor: str
    text_preview: str
    is_merge_origin: bool
    is_spanned: bool
    span_width: int
    span_height: int


@dataclass(frozen=True)
class TableParagraphFormatRecord:
    """表格单元格段落级缩进事实。"""

    slide_number: int
    owner_shape_id: int
    owner_shape_name: str
    table_index_on_slide: int
    row: int
    col: int
    paragraph: int
    text_preview: str
    paragraph_level: int
    paragraph_margin_left_pt: float
    paragraph_margin_right_pt: float
    paragraph_indent_pt: float
    nonzero_indents: dict[str, float]


@dataclass(frozen=True)
class TextLayoutMetrics:
    """统一的文本布局估算结果。"""

    estimated_bounds_pt: RectPt
    line_count_estimated: int
    estimated_height_pt: float
    single_line_width_pt: float
    effective_inner_width_pt: float
    width_pressure_ratio: float
    max_line_width_pt: float
    last_line_width_pt: float
    bottom_gap_pt: float
    right_gap_pt: float
    overflow_area_pt2: float
    overflow_ratio: float


@dataclass(frozen=True)
class QualityIssue:
    """统一的问题记录。"""

    severity: str
    issue_type: str
    message: str
    slide_number: int | None = None
    shape_id: int | None = None
    source_kind: str | None = None
    details: dict | None = None
    suggested_fix: str | None = None


def emu_to_pt(value) -> float:
    """把 EMU 转换成 pt。"""
    if value is None:
        return 0.0
    return round(float(value) / EMU_PER_PT, 2)


def rect_from_shape(shape) -> RectPt:
    """从 shape 构造矩形。"""
    return RectPt(
        left=emu_to_pt(shape.left),
        top=emu_to_pt(shape.top),
        width=emu_to_pt(shape.width),
        height=emu_to_pt(shape.height),
    )


def normalize_text(text: str) -> str:
    """压缩空白，返回单行文本。"""
    return re.sub(r"\s+", " ", text or "").strip()


def shorten_text(text: str, limit: int = 120) -> str:
    """截断长文本。"""
    if len(text) <= limit:
        return text
    return text[: limit - 1] + "…"


def shape_type_name(shape) -> str:
    """返回稳定的 shape type 标签。"""
    shape_type = getattr(shape, "shape_type", None)
    if shape_type is None:
        return "UNKNOWN"
    return getattr(shape_type, "name", str(shape_type))


def iter_shape_records(slide_number: int, shapes, parent_group_id: int | None = None) -> Iterable[ShapeRecord]:
    """递归遍历 slide shapes，输出统一记录。"""
    for z_order, shape in enumerate(shapes, start=1):
        record = ShapeRecord(
            slide_number=slide_number,
            shape_id=getattr(shape, "shape_id", z_order),
            shape_name=getattr(shape, "name", ""),
            shape_type=shape_type_name(shape),
            z_order=z_order,
            parent_group_id=parent_group_id,
            rect_pt=rect_from_shape(shape),
            has_text=getattr(shape, "has_text_frame", False),
            has_table=getattr(shape, "has_table", False),
            has_chart=getattr(shape, "has_chart", False),
            is_picture=getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE,
            source_kind="group"
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP
            else ("picture" if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE else "shape"),
        )
        yield record
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shape_records(slide_number, shape.shapes, parent_group_id=record.shape_id)


def collect_shape_inventory(prs: Presentation) -> list[ShapeRecord]:
    """收集整份 PPT 的 shape inventory。"""
    inventory: list[ShapeRecord] = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        inventory.extend(iter_shape_records(slide_number, slide.shapes))
    return inventory


def pick_paragraph_font_size(paragraph, fallback: float = DEFAULT_BODY_FONT_PT) -> float:
    """从 paragraph / runs 中选一个更稳的字号。"""
    sizes: list[float] = []
    paragraph_size = getattr(paragraph.font, "size", None)
    if paragraph_size is not None:
        sizes.append(emu_to_pt(paragraph_size))
    for run in paragraph.runs:
        run_size = getattr(run.font, "size", None)
        if run_size is not None:
            sizes.append(emu_to_pt(run_size))
    if not sizes:
        return fallback
    return max(size for size in sizes if size > 0)


def collect_font_size_occurrences(prs: Presentation) -> list[FontSizeOccurrence]:
    """收集 slide 可见文本中能够解析到的逐 run 字号。"""

    occurrences: list[FontSizeOccurrence] = []

    def collect_text_frame(slide_number: int, shape, text_frame, source_kind: str) -> None:
        shape_id = getattr(shape, "shape_id", -1)
        shape_name = getattr(shape, "name", "")
        for paragraph in text_frame.paragraphs:
            paragraph_size = getattr(paragraph.font, "size", None)
            paragraph_size_pt = emu_to_pt(paragraph_size) if paragraph_size is not None else None
            for run in paragraph.runs:
                text = normalize_text(run.text)
                if not text:
                    continue
                run_size = getattr(run.font, "size", None)
                if run_size is not None:
                    font_size_pt = emu_to_pt(run_size)
                    setting_source = "run"
                elif paragraph_size_pt is not None:
                    font_size_pt = paragraph_size_pt
                    setting_source = "paragraph"
                else:
                    continue
                if font_size_pt <= 0:
                    continue
                occurrences.append(
                    FontSizeOccurrence(
                        slide_number=slide_number,
                        owner_shape_id=shape_id,
                        owner_shape_name=shape_name,
                        source_kind=source_kind,
                        text=shorten_text(text, 80),
                        font_size_pt=font_size_pt,
                        setting_source=setting_source,
                    )
                )

    def visit_shapes(slide_number: int, shapes) -> None:
        for shape in shapes:
            if getattr(shape, "has_text_frame", False):
                collect_text_frame(slide_number, shape, shape.text_frame, "shape_text")
            if getattr(shape, "has_table", False):
                for row_index, row in enumerate(shape.table.rows):
                    for col_index, cell in enumerate(row.cells):
                        collect_text_frame(
                            slide_number,
                            shape,
                            cell.text_frame,
                            f"table_cell[{row_index},{col_index}]",
                        )
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                visit_shapes(slide_number, shape.shapes)

    for slide_number, slide in enumerate(prs.slides, start=1):
        visit_shapes(slide_number, slide.shapes)
    return occurrences


def text_frame_inner_rect(shape) -> RectPt:
    """估算 text frame 的可用内容区。"""
    rect = rect_from_shape(shape)
    text_frame = shape.text_frame
    margin_left = emu_to_pt(getattr(text_frame, "margin_left", 0))
    margin_right = emu_to_pt(getattr(text_frame, "margin_right", 0))
    margin_top = emu_to_pt(getattr(text_frame, "margin_top", 0))
    margin_bottom = emu_to_pt(getattr(text_frame, "margin_bottom", 0))
    return RectPt(
        left=rect.left + margin_left,
        top=rect.top + margin_top,
        width=max(0.0, rect.width - margin_left - margin_right),
        height=max(0.0, rect.height - margin_top - margin_bottom),
    )


def estimate_text_width_units(text: str) -> float:
    """按字符类别粗略估计一行文本宽度单位。"""
    units = 0.0
    for char in text:
        if char.isspace():
            units += 0.35
        elif "\u4e00" <= char <= "\u9fff":
            units += 1.0
        elif char.isdigit():
            units += 0.56
        elif "A" <= char <= "Z":
            units += 0.66
        elif "a" <= char <= "z":
            units += 0.56
        else:
            units += 0.42
    return units


def estimate_line_count(text: str, font_size_pt: float, inner_width_pt: float) -> int:
    """按启发式估算换行后的行数。"""
    if not text:
        return 0
    if inner_width_pt <= 0:
        return 1

    char_unit_width_pt = font_size_pt * 0.62
    max_units = max(1.0, inner_width_pt / char_unit_width_pt)

    line_count = 0
    for raw_line in text.split("\n"):
        normalized = normalize_text(raw_line)
        if not normalized:
            line_count += 1
            continue
        units = estimate_text_width_units(normalized)
        line_count += max(1, math.ceil(units / max_units))
    return line_count


def estimate_text_layout_metrics(
    inner_rect: RectPt,
    text: str,
    font_size_pt: float,
    paragraph_count: int,
) -> TextLayoutMetrics:
    """估算文本布局和边界风险指标。"""
    if not text:
        empty_rect = RectPt(inner_rect.left, inner_rect.top, 0.0, 0.0)
        return TextLayoutMetrics(
            estimated_bounds_pt=empty_rect,
            line_count_estimated=0,
            estimated_height_pt=0.0,
            single_line_width_pt=0.0,
            effective_inner_width_pt=inner_rect.width,
            width_pressure_ratio=0.0,
            max_line_width_pt=0.0,
            last_line_width_pt=0.0,
            bottom_gap_pt=inner_rect.height,
            right_gap_pt=inner_rect.width,
            overflow_area_pt2=0.0,
            overflow_ratio=0.0,
        )

    char_unit_width_pt = font_size_pt * 0.62
    max_units = max(1.0, inner_rect.width / char_unit_width_pt) if inner_rect.width > 0 else 1.0
    effective_inner_width_pt = max(1.0, inner_rect.width - font_size_pt * 0.7)
    single_line_units = max(estimate_text_width_units(normalize_text(text)), 0.0)
    single_line_width_pt = round(single_line_units * char_unit_width_pt, 2)
    width_pressure_ratio = round(single_line_width_pt / effective_inner_width_pt, 4)
    displayed_line_units: list[float] = []
    for raw_line in text.split("\n"):
        normalized = normalize_text(raw_line)
        if not normalized:
            displayed_line_units.append(0.0)
            continue
        units = estimate_text_width_units(normalized)
        segment_count = max(1, math.ceil(units / max_units))
        remaining_units = units
        for segment_index in range(segment_count):
            if segment_index < segment_count - 1:
                displayed_line_units.append(max_units)
                remaining_units -= max_units
            else:
                displayed_line_units.append(max(0.0, remaining_units) or max_units)

    line_count = len(displayed_line_units)
    line_height_pt = font_size_pt * 1.08
    paragraph_gap_pt = max(0.0, (paragraph_count - 1) * font_size_pt * 0.18)
    estimated_height_pt = round(line_count * line_height_pt + paragraph_gap_pt, 2)
    max_line_width_pt = round(min(inner_rect.width, max(displayed_line_units, default=0.0) * char_unit_width_pt), 2)
    last_line_width_pt = round(min(inner_rect.width, displayed_line_units[-1] * char_unit_width_pt), 2) if displayed_line_units else 0.0
    bottom_gap_pt = round(inner_rect.height - estimated_height_pt, 2)
    right_gap_pt = round(inner_rect.width - last_line_width_pt, 2)
    overflow_height_pt = max(0.0, estimated_height_pt - inner_rect.height)
    overflow_width_pt = max(0.0, max_line_width_pt - inner_rect.width)
    overflow_area_pt2 = round(
        overflow_height_pt * max_line_width_pt
        + overflow_width_pt * estimated_height_pt
        - overflow_height_pt * overflow_width_pt,
        2,
    )
    estimated_text_area = max(1.0, max_line_width_pt * estimated_height_pt)
    overflow_ratio = round(overflow_area_pt2 / estimated_text_area, 4)

    estimated_bounds = RectPt(
        left=inner_rect.left,
        top=inner_rect.top,
        width=max_line_width_pt,
        height=estimated_height_pt,
    )
    return TextLayoutMetrics(
        estimated_bounds_pt=estimated_bounds,
        line_count_estimated=line_count,
        estimated_height_pt=estimated_height_pt,
        single_line_width_pt=single_line_width_pt,
        effective_inner_width_pt=effective_inner_width_pt,
        width_pressure_ratio=width_pressure_ratio,
        max_line_width_pt=max_line_width_pt,
        last_line_width_pt=last_line_width_pt,
        bottom_gap_pt=bottom_gap_pt,
        right_gap_pt=right_gap_pt,
        overflow_area_pt2=overflow_area_pt2,
        overflow_ratio=overflow_ratio,
    )


def estimate_text_bounds(inner_rect: RectPt, text: str, font_size_pt: float, paragraph_count: int) -> RectPt:
    """为兼容旧调用保留的文本边界估算接口。"""
    return estimate_text_layout_metrics(inner_rect, text, font_size_pt, paragraph_count).estimated_bounds_pt


def collect_text_items(prs: Presentation) -> list[TextItemRecord]:
    """收集 shape / table cell 级文本对象。"""
    items: list[TextItemRecord] = []
    inventory_index = {
        (record.slide_number, record.shape_id): record for record in collect_shape_inventory(prs)
    }

    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            owner = inventory_index.get((slide_number, getattr(shape, "shape_id", -1)))

            if getattr(shape, "has_text_frame", False):
                paragraphs = [paragraph for paragraph in shape.text_frame.paragraphs if normalize_text(paragraph.text)]
                if paragraphs:
                    text = "\n".join(normalize_text(paragraph.text) for paragraph in paragraphs)
                    font_size_pt = max(pick_paragraph_font_size(paragraph) for paragraph in paragraphs)
                    inner_rect = text_frame_inner_rect(shape)
                    items.append(
                        TextItemRecord(
                            slide_number=slide_number,
                            owner_shape_id=owner.shape_id if owner else getattr(shape, "shape_id", -1),
                            owner_shape_name=owner.shape_name if owner else getattr(shape, "name", ""),
                            owner_shape_type=owner.shape_type if owner else shape_type_name(shape),
                            z_order=owner.z_order if owner else 0,
                            source_kind="shape_text",
                            text=shorten_text(text),
                            rect_pt=rect_from_shape(shape),
                            inner_rect_pt=inner_rect,
                            font_size_pt=font_size_pt,
                            paragraph_count=len(paragraphs),
                        )
                    )

            if getattr(shape, "has_table", False):
                base_rect = rect_from_shape(shape)
                current_top = base_rect.top
                for row_index, row in enumerate(shape.table.rows):
                    row_height_pt = emu_to_pt(row.height)
                    current_left = base_rect.left
                    for col_index, cell in enumerate(row.cells):
                        col_width_pt = emu_to_pt(shape.table.columns[col_index].width)
                        text = normalize_text(cell.text)
                        if text:
                            inner_rect = RectPt(
                                left=current_left + 4.0,
                                top=current_top + 3.0,
                                width=max(0.0, col_width_pt - 8.0),
                                height=max(0.0, row_height_pt - 6.0),
                            )
                            paragraphs = [paragraph for paragraph in cell.text_frame.paragraphs if normalize_text(paragraph.text)]
                            font_size_pt = max(
                                (pick_paragraph_font_size(paragraph) for paragraph in paragraphs),
                                default=DEFAULT_BODY_FONT_PT,
                            )
                            items.append(
                                TextItemRecord(
                                    slide_number=slide_number,
                                    owner_shape_id=owner.shape_id if owner else getattr(shape, "shape_id", -1),
                                    owner_shape_name=f"{getattr(shape, 'name', '')}[{row_index},{col_index}]",
                                    owner_shape_type=owner.shape_type if owner else shape_type_name(shape),
                                    z_order=owner.z_order if owner else 0,
                                    source_kind="table_cell",
                                    text=shorten_text(text),
                                    rect_pt=RectPt(current_left, current_top, col_width_pt, row_height_pt),
                                    inner_rect_pt=inner_rect,
                                    font_size_pt=font_size_pt,
                                    paragraph_count=max(1, len(paragraphs)),
                                )
                            )
                        current_left += col_width_pt
                    current_top += row_height_pt
    return items


def collect_table_cell_format_records(prs: Presentation) -> list[TableCellFormatRecord]:
    """收集原生 PPT 表格单元格的垂直对齐事实。"""

    records: list[TableCellFormatRecord] = []

    def visit_shapes(slide_number: int, shapes, table_counter: list[int]) -> None:
        for shape in shapes:
            if getattr(shape, "has_table", False):
                table_counter[0] += 1
                records.extend(iter_table_cell_format_records(slide_number, shape, table_counter[0]))
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                visit_shapes(slide_number, shape.shapes, table_counter)

    for slide_number, slide in enumerate(prs.slides, start=1):
        visit_shapes(slide_number, slide.shapes, [0])
    return records


def collect_table_paragraph_format_records(prs: Presentation) -> list[TableParagraphFormatRecord]:
    """收集原生 PPT 表格单元格段落的缩进事实。"""

    records: list[TableParagraphFormatRecord] = []

    def visit_shapes(slide_number: int, shapes, table_counter: list[int]) -> None:
        for shape in shapes:
            if getattr(shape, "has_table", False):
                table_counter[0] += 1
                records.extend(iter_table_paragraph_format_records(slide_number, shape, table_counter[0]))
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                visit_shapes(slide_number, shape.shapes, table_counter)

    for slide_number, slide in enumerate(prs.slides, start=1):
        visit_shapes(slide_number, slide.shapes, [0])
    return records


def iter_table_cell_format_records(slide_number: int, shape, table_index_on_slide: int) -> Iterable[TableCellFormatRecord]:
    """遍历单个 table shape 的非空、可见单元格格式事实。"""

    shape_id = getattr(shape, "shape_id", -1)
    shape_name = getattr(shape, "name", "")
    for row_index, row in enumerate(shape.table.rows, start=1):
        for col_index, cell in enumerate(row.cells, start=1):
            text = normalize_text(cell.text)
            if not text or getattr(cell, "is_spanned", False):
                continue
            records = TableCellFormatRecord(
                slide_number=slide_number,
                owner_shape_id=shape_id,
                owner_shape_name=shape_name,
                table_index_on_slide=table_index_on_slide,
                row=row_index,
                col=col_index,
                vertical_anchor=vertical_anchor_name(cell.vertical_anchor),
                text_preview=shorten_text(text, 80),
                is_merge_origin=getattr(cell, "is_merge_origin", False),
                is_spanned=getattr(cell, "is_spanned", False),
                span_width=getattr(cell, "span_width", 1),
                span_height=getattr(cell, "span_height", 1),
            )
            yield records


def iter_table_paragraph_format_records(slide_number: int, shape, table_index_on_slide: int) -> Iterable[TableParagraphFormatRecord]:
    """遍历单个 table shape 的非空单元格段落缩进事实。"""

    shape_id = getattr(shape, "shape_id", -1)
    shape_name = getattr(shape, "name", "")
    for row_index, row in enumerate(shape.table.rows, start=1):
        for col_index, cell in enumerate(row.cells, start=1):
            if getattr(cell, "is_spanned", False):
                continue
            for paragraph_index, paragraph in enumerate(cell.text_frame.paragraphs, start=1):
                text = normalize_text(paragraph.text)
                if not text:
                    continue
                indent_values = paragraph_indent_values(paragraph)
                nonzero_indents = {
                    name: value
                    for name, value in indent_values.items()
                    if not is_zero_table_indent(value)
                }
                yield TableParagraphFormatRecord(
                    slide_number=slide_number,
                    owner_shape_id=shape_id,
                    owner_shape_name=shape_name,
                    table_index_on_slide=table_index_on_slide,
                    row=row_index,
                    col=col_index,
                    paragraph=paragraph_index,
                    text_preview=shorten_text(text, 80),
                    nonzero_indents=nonzero_indents,
                    **indent_values,
                )


def paragraph_indent_values(paragraph) -> dict[str, float | int]:
    """读取 PPT 段落层级和 OOXML 缩进属性，未设置时按零处理。"""

    return {
        "paragraph_level": paragraph_level(paragraph),
        "paragraph_margin_left_pt": paragraph_indent_attr_pt(paragraph, "marL"),
        "paragraph_margin_right_pt": paragraph_indent_attr_pt(paragraph, "marR"),
        "paragraph_indent_pt": paragraph_indent_attr_pt(paragraph, "indent"),
    }


def paragraph_level(paragraph) -> int:
    """读取 `a:pPr/@lvl`，缺省时返回顶层级别 `0`。"""

    p_pr = paragraph._p.pPr
    if p_pr is None:
        return 0
    raw_value = p_pr.get("lvl")
    if raw_value is None:
        return 0
    try:
        return int(raw_value)
    except ValueError as exc:
        raise ValueError(f"非法 PPT 段落级别值: lvl={raw_value!r}") from exc


def paragraph_indent_attr_pt(paragraph, attribute_name: str) -> float:
    """读取 `a:pPr` 段落缩进属性并转换为 pt，缺省时返回零。"""

    p_pr = paragraph._p.pPr
    if p_pr is None:
        return 0.0
    raw_value = p_pr.get(attribute_name)
    if raw_value is None:
        return 0.0
    try:
        return emu_to_pt(int(raw_value))
    except ValueError as exc:
        raise ValueError(f"非法 PPT 段落缩进值: {attribute_name}={raw_value!r}") from exc


def vertical_anchor_name(anchor) -> str:
    """把 `python-pptx` 的垂直锚点枚举转成稳定小写标签。"""

    if anchor is None:
        return "unset"
    name = getattr(anchor, "name", None)
    return str(name).lower() if name else str(anchor).lower()


def is_zero_table_indent(value: float) -> bool:
    """判断表格段落缩进是否可视为零。"""

    return abs(float(value)) <= TABLE_INDENT_TOLERANCE_PT


def rect_intersection_area(rect_a: RectPt, rect_b: RectPt) -> float:
    """返回两个矩形的相交面积。"""
    overlap_width = max(0.0, min(rect_a.right, rect_b.right) - max(rect_a.left, rect_b.left))
    overlap_height = max(0.0, min(rect_a.bottom, rect_b.bottom) - max(rect_a.top, rect_b.top))
    return round(overlap_width * overlap_height, 2)


def shape_can_occlude(record: ShapeRecord) -> bool:
    """判断 shape 是否可能遮挡文字。"""
    if record.shape_type in {"LINE", "STRAIGHT_CONNECTOR_1", "ELBOW_CONNECTOR_2", "CURVED_CONNECTOR_3"}:
        return False
    return record.rect_pt.width > 0 and record.rect_pt.height > 0


def issue_summary(issues: list[QualityIssue]) -> dict[str, int]:
    """统计问题数量。"""
    summary = {"error": 0, "warning": 0, "not_checked": 0}
    for issue in issues:
        summary.setdefault(issue.severity, 0)
        summary[issue.severity] += 1
    return summary


def dump_json(path: Path, payload: dict) -> None:
    """写出 JSON。"""
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def report_timestamp() -> str:
    """返回适合写入文件名的时间戳。"""
    return datetime.now().strftime("%Y%m%d_%H%M%S")


def resolve_gate_report_paths(
    *,
    gate_name: str,
    workspace_dir: Path | None,
    json_out: Path | None,
    md_out: Path | None,
) -> tuple[Path | None, Path | None, str | None]:
    """为 quality gate 生成标准报告路径。"""
    if workspace_dir is None:
        return json_out, md_out, None

    timestamp = report_timestamp()
    history_dir = workspace_dir.resolve() / "validation" / gate_name / "history"
    resolved_json = json_out or (history_dir / f"{gate_name}_{timestamp}.json")
    resolved_md = md_out or (history_dir / f"{gate_name}_{timestamp}.md")
    return resolved_json, resolved_md, timestamp


def render_issue_markdown(title: str, payload: dict) -> str:
    """把问题报告渲染成高效摘要版 Markdown。"""
    summary = payload.get("summary", {})
    issues = payload.get("issues", [])
    lines = [
        f"# {title}",
        "",
        f"- 输入文件：`{payload.get('pptx', '')}`",
        f"- 错误数：`{summary.get('error', 0)}`",
        f"- 警告数：`{summary.get('warning', 0)}`",
        f"- 未检查数：`{summary.get('not_checked', 0)}`",
        "",
    ]

    if not issues:
        lines.append("## 结果")
        lines.append("")
        lines.append("未发现问题。")
        lines.append("")
        return "\n".join(lines)

    grouped: dict[tuple[str, str, str, str], list[dict]] = {}
    for issue in issues:
        grouped.setdefault(
            (
                issue["severity"],
                issue["issue_type"],
                issue["message"],
                issue.get("suggested_fix") or "",
            ),
            [],
        ).append(issue)

    lines.extend(["## 摘要", ""])
    for (severity, issue_type, _, _), bucket in sorted(grouped.items(), key=lambda item: (item[0][0], item[0][1])):
        lines.append(f"- `{severity}` / `{issue_type}`: {len(bucket)}")
    lines.append("")

    def _severity_sort_key(name: str) -> int:
        order = {"error": 0, "warning": 1, "not_checked": 2}
        return order.get(name, 99)

    def _fmt_number(value) -> str | None:
        if value is None:
            return None
        if isinstance(value, float):
            text = f"{value:.4f}".rstrip("0").rstrip(".")
            return text
        return str(value)

    def _location_key(issue: dict) -> tuple:
        details = issue.get("details") or {}
        return (
            issue.get("slide_number"),
            issue.get("shape_id"),
            details.get("table_index_on_slide"),
            details.get("row"),
            details.get("col"),
            details.get("paragraph"),
            details.get("target_shape_id"),
            details.get("occluding_shape_id"),
        )

    def _location_brief(location_issues: list[dict]) -> str:
        exemplar = location_issues[0]
        details_list = [issue.get("details") or {} for issue in location_issues]
        parts: list[str] = []
        if exemplar.get("slide_number") is not None:
            parts.append(f"slide {exemplar['slide_number']}")
        if exemplar.get("shape_id") is not None:
            parts.append(f"shape {exemplar['shape_id']}")

        table_positions = {
            (
                details.get("table_index_on_slide"),
                details.get("row"),
                details.get("col"),
                details.get("paragraph"),
            )
            for details in details_list
            if details.get("row") is not None and details.get("col") is not None
        }
        if table_positions:
            formatted_positions = []
            for table_index, row, col, paragraph in sorted(table_positions):
                item = f"table={table_index} row={row} col={col}"
                if paragraph is not None:
                    item += f" paragraph={paragraph}"
                formatted_positions.append(item)
            parts.append("cells=" + "; ".join(formatted_positions[:3]))
            if len(formatted_positions) > 3:
                parts.append(f"cell_occurrences={len(formatted_positions)}")

        target_shape_ids = {details.get("target_shape_id") for details in details_list if details.get("target_shape_id") is not None}
        occluding_shape_ids = {details.get("occluding_shape_id") for details in details_list if details.get("occluding_shape_id") is not None}
        if target_shape_ids:
            parts.append("target_shape=" + ",".join(str(value) for value in sorted(target_shape_ids)))
        if occluding_shape_ids:
            parts.append("occluder=" + ",".join(str(value) for value in sorted(occluding_shape_ids)))
        if len(location_issues) > 1:
            parts.append(f"occurrences={len(location_issues)}")

        font_sizes: set[float] = set()
        nearest_sizes: set[float] = set()
        for details in details_list:
            if details.get("font_size_pt") is not None:
                font_sizes.add(details["font_size_pt"])
            font_sizes.update(details.get("font_sizes_pt") or [])
            if details.get("nearest_half_point_pt") is not None:
                nearest_sizes.add(details["nearest_half_point_pt"])
        if font_sizes:
            parts.append("font_sizes_pt=" + ",".join(_fmt_number(value) or "" for value in sorted(font_sizes)))
        if nearest_sizes:
            parts.append("nearest_half_point_pt=" + ",".join(_fmt_number(value) or "" for value in sorted(nearest_sizes)))

        def _pick_metric(key: str, reducer=max):
            values = [details.get(key) for details in details_list if details.get(key) is not None]
            if not values:
                return None
            return reducer(values)

        overflow_ratio = _pick_metric("overflow_ratio", max)
        bottom_gap = _pick_metric("bottom_gap_pt", min)
        right_gap = _pick_metric("right_gap_pt", min)
        overlap_ratio = _pick_metric("overlap_ratio", max)
        if overflow_ratio is not None:
            parts.append(f"overflow_ratio={_fmt_number(overflow_ratio)}")
        if bottom_gap is not None:
            parts.append(f"bottom_gap_pt={_fmt_number(bottom_gap)}")
        if right_gap is not None:
            parts.append(f"right_gap_pt={_fmt_number(right_gap)}")
        if overlap_ratio is not None:
            parts.append(f"overlap_ratio={_fmt_number(overlap_ratio)}")

        strip_analyses = [details.get("strip_analysis") for details in details_list if details.get("strip_analysis")]
        if strip_analyses:
            dark_ratio = max(item.get("dark_ratio", 0) for item in strip_analyses)
            component = max(item.get("max_component_area", 0) for item in strip_analyses)
            parts.append(f"dark_ratio={_fmt_number(dark_ratio)}")
            parts.append(f"component={component}")

        return " | ".join(parts)

    severity_buckets: dict[str, list[tuple[tuple[str, str, str, str], list[dict]]]] = {}
    for key, bucket in grouped.items():
        severity_buckets.setdefault(key[0], []).append((key, bucket))

    lines.extend(["## 问题清单", ""])
    for severity in sorted(severity_buckets, key=_severity_sort_key):
        lines.append(f"### `{severity}`")
        lines.append("")
        for (severity_name, issue_type, message, suggested_fix), bucket in sorted(
            severity_buckets[severity], key=lambda item: item[0][1]
        ):
            lines.append(f"#### `{issue_type}`")
            lines.append("")
            lines.append(message)
            lines.append("")
            if suggested_fix:
                lines.append(f"建议：{suggested_fix}")
                lines.append("")

            by_location: dict[tuple, list[dict]] = {}
            for issue in bucket:
                by_location.setdefault(_location_key(issue), []).append(issue)
            lines.append("出现位置：")
            sorted_locations = sorted(
                by_location.items(),
                key=lambda item: (
                    item[0][0] if item[0][0] is not None else 0,
                    item[0][1] if item[0][1] is not None else 0,
                ),
            )
            location_limit = (
                FONT_SIZE_LOCATION_EXAMPLE_LIMIT
                if issue_type.startswith("font_size_") or "_font_" in issue_type
                else len(sorted_locations)
            )
            for _, location_issues in sorted_locations[:location_limit]:
                lines.append(f"- {_location_brief(location_issues)}")
            omitted_count = len(sorted_locations) - location_limit
            if omitted_count > 0:
                lines.append(f"- 其余 {omitted_count} 个位置见 JSON 报告")
            lines.append("")
    return "\n".join(lines)


def write_issue_bundle(
    *,
    title: str,
    pptx_path: Path,
    issues: list[QualityIssue],
    json_out: Path | None,
    md_out: Path | None,
    extra_payload: dict | None = None,
    generated_at: str | None = None,
) -> dict:
    """统一写出 JSON / Markdown 报告。"""
    payload = {
        "pptx": str(pptx_path.resolve()),
        "generated_at": generated_at,
        "summary": issue_summary(issues),
        "issues": [asdict(issue) for issue in issues],
    }
    if extra_payload:
        payload.update(extra_payload)

    if json_out:
        dump_json(json_out, payload)
    if md_out:
        md_out.parent.mkdir(parents=True, exist_ok=True)
        md_out.write_text(render_issue_markdown(title, payload), encoding="utf-8")

    gate_name = infer_gate_name_from_title(title)
    observations = quality_issues_to_observations(
        issues,
        skill="ppt-polished-deck-collab",
        gate=gate_name,
        artifact_path=pptx_path,
        detector=detector_for_gate(gate_name),
    )
    reminder = build_agent_reminder(
        observations,
        source={"skill": "ppt-polished-deck-collab", "gate": gate_name},
        artifact={"path": str(pptx_path.resolve())},
        full_report_ref=str(json_out) if json_out else None,
        target_milestone="final",
    )
    agent_json_out = sidecar_path(json_out, ".json")
    agent_md_out = sidecar_path(md_out, ".md")
    write_agent_reminder(json_path=agent_json_out, md_path=agent_md_out, reminder=reminder)
    payload["agent_reminder"] = {
        "json": str(agent_json_out) if agent_json_out else None,
        "markdown": str(agent_md_out) if agent_md_out else None,
        "decision": reminder["decision"],
    }
    return payload


def infer_gate_name_from_title(title: str) -> str:
    """从报告标题推断 gate 名称。"""

    normalized = title.lower()
    if "package" in normalized:
        return "package_preflight"
    if "structure" in normalized:
        return "structure_precheck"
    if "render" in normalized:
        return "render_review"
    return "quality_gate"


def detector_for_gate(gate_name: str) -> dict:
    """返回 gate 对应的 detector 元数据。"""

    mapping = {
        "package_preflight": {
            "id": "ppt.package_preflight",
            "layer": "package",
            "method": "ooxml_package_analysis",
            "version": "1.0",
        },
        "structure_precheck": {
            "id": "ppt.structure_precheck",
            "layer": "structure",
            "method": "pptx_geometry_analysis",
            "version": "1.0",
        },
        "render_review": {
            "id": "ppt.render_review",
            "layer": "raster",
            "method": "raster_edge_detection",
            "version": "1.0",
        },
    }
    return mapping.get(
        gate_name,
        {"id": f"ppt.{gate_name}", "layer": "unknown", "method": "unknown", "version": "1.0"},
    )
